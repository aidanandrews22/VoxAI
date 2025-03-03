"""
Presentation processor module for extracting text from PowerPoint presentations.
"""
import io
import os
import tempfile
from typing import Dict, Any, List, Tuple

from PIL import Image
from pptx import Presentation

from app.core.logging import logger
from app.services.file_processors import FileProcessor


class PresentationProcessor(FileProcessor):
    """
    Base processor for presentation files.
    """
    
    async def process(self, file_content: bytes, file_path: str) -> str:
        """
        Process presentation content and extract text.
        
        Args:
            file_content: Raw presentation file content bytes
            file_path: Path to the presentation file
            
        Returns:
            str: Extracted text content
        """
        raise NotImplementedError("Subclasses must implement this method")
    
    async def get_metadata(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from presentation file.
        
        Args:
            file_content: Raw presentation file content bytes
            file_path: Path to the presentation file
            
        Returns:
            Dict[str, Any]: Extracted metadata
        """
        return {
            'file_size': len(file_content),
            'file_extension': os.path.splitext(file_path)[1],
        }

    async def _extract_text_and_description_from_image(self, image: Image.Image) -> str:
        """
        Extract text and generate description from an image using Google Gemini.
        
        Args:
            image: PIL Image object
            
        Returns:
            str: Extracted text and generated description
        """
        try:
            from google import genai
            import os
            
            # Initialize the Gemini client
            client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
            
            # Convert PIL Image to bytes for the API
            with io.BytesIO() as output:
                image.save(output, format="PNG")
                image_bytes = output.getvalue()
            
            # Create the request parts
            from google.genai import types
            image_part = types.Part.from_bytes(data=image_bytes, mime_type="image/png")
            
            combined_prompt = """
                1. Extract all text visible in this image. If no text is visible, respond with 'No text detected.'
                2. Provide a detailed description of what's in this image.
            """

            response = client.models.generate_content(
                model="gemini-1.5-flash-8b",
                contents=[combined_prompt, image_part]
            )
            
            return response.text
        except Exception as e:
            logger.error(f"Error extracting text and description from image: {e}")
            return "Error processing image content."


class PowerPointProcessor(PresentationProcessor):
    """
    Processor for PowerPoint files.
    Handles: application/vnd.openxmlformats-officedocument.presentationml.presentation
    """
    
    async def process(self, file_content: bytes, file_path: str) -> str:
        """
        Process PowerPoint content and extract text.
        
        Args:
            file_content: Raw PowerPoint file content bytes
            file_path: Path to the PowerPoint file
            
        Returns:
            str: Extracted text content
        """
        try:
            # Write content to a temporary file because python-pptx requires a file path
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(file_content)
                temp_file.flush()
                temp_path = temp_file.name
            
            # Load the presentation
            ppt = Presentation(temp_path)
            
            # Extract text from all slides
            full_text = []
            
            for slide_num, slide in enumerate(ppt.slides, start=1):
                slide_text = []
                slide_text.append(f"Slide {slide_num}")
                slide_text.append("=" * (len(f"Slide {slide_num}")))
                
                # Extract text from shapes with position information
                shape_contents = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Get position information
                        top = shape.top if hasattr(shape, "top") else 0
                        shape_contents.append((shape.text, top))
                
                # Extract images with position information
                image_list = await self._extract_images_with_positions(slide)
                image_contents = []
                
                for img, left, top in image_list:
                    img_content = await self._extract_text_and_description_from_image(img)
                    if img_content:
                        image_contents.append((f"[IMAGE CONTENT START]\n{img_content}\n[IMAGE CONTENT END]", top))
                
                # Combine all content and sort by vertical position
                all_contents = shape_contents + image_contents
                all_contents.sort(key=lambda x: x[1])  # Sort by top position
                
                # Extract just the content text after sorting
                sorted_contents = [content for content, _ in all_contents]
                
                # Determine slide content type and organize accordingly
                if not sorted_contents:
                    # Empty slide
                    slide_text.append("[EMPTY SLIDE]")
                else:
                    # Add all content in position order
                    slide_text.extend(sorted_contents)
                
                # Create a text summary for this slide
                if len(slide_text) > 2:  # If we have more than just the header
                    full_text.append("\n".join(slide_text))
            
            # Clean up temporary file
            os.unlink(temp_path)
            
            return "\n\n".join(full_text)
        except Exception as e:
            logger.error(f"Error processing PowerPoint file: {e}")
            raise
    
    async def get_metadata(self, file_content: bytes, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from PowerPoint file.
        
        Args:
            file_content: Raw PowerPoint file content bytes
            file_path: Path to the PowerPoint file
            
        Returns:
            Dict[str, Any]: Extracted metadata
        """
        base_metadata = await super().get_metadata(file_content, file_path)
        
        try:
            # Write content to a temporary file because python-pptx requires a file path
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(file_content)
                temp_file.flush()
                temp_path = temp_file.name
            
            # Load the presentation
            ppt = Presentation(temp_path)
            
            # Extract presentation metadata
            core_props = ppt.core_properties
            metadata = {
                'slide_count': len(ppt.slides),
                'title': core_props.title,
                'author': core_props.author,
                'subject': core_props.subject,
                'keywords': core_props.keywords,
                'created': core_props.created.isoformat() if core_props.created else None,
                'modified': core_props.modified.isoformat() if core_props.modified else None,
                'last_modified_by': core_props.last_modified_by,
            }
            
            # Count shapes by type
            shape_counts = {}
            total_image_count = 0
            for slide in ppt.slides:
                for shape in slide.shapes:
                    shape_type = shape.shape_type
                    shape_counts[shape_type] = shape_counts.get(shape_type, 0) + 1
                    
                    # Count images specifically
                    if hasattr(shape, 'image'):
                        total_image_count += 1
            
            if shape_counts:
                metadata['shape_counts'] = shape_counts
            
            metadata['total_image_count'] = total_image_count
            
            # Extract slide information
            slide_info = []
            for i, slide in enumerate(ppt.slides):
                text_length = 0
                image_count = 0
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_length += len(shape.text)
                    if hasattr(shape, 'image'):
                        image_count += 1
                
                slide_info.append({
                    'slide_number': i + 1,
                    'shape_count': len(slide.shapes),
                    'text_length': text_length,
                    'image_count': image_count
                })
            
            if slide_info:
                metadata['slides'] = slide_info
            
            # Remove empty metadata fields
            metadata = {k: v for k, v in metadata.items() if v}
            
            # Add to base metadata
            base_metadata.update(metadata)
            
            # Clean up temporary file
            os.unlink(temp_path)
            
            return base_metadata
        except Exception as e:
            logger.error(f"Error extracting PowerPoint metadata: {e}")
            return base_metadata
    
    async def _extract_images_from_slide(self, slide) -> List[Image.Image]:
        """
        Extract images from a PowerPoint slide.
        
        Args:
            slide: python-pptx slide object
            
        Returns:
            List[Image.Image]: List of extracted images as PIL Image objects
        """
        images = []
        
        for shape in slide.shapes:
            try:
                # Check if shape is a picture
                if hasattr(shape, 'image'):
                    try:
                        # Extract the image data
                        image_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        
                        # Only process images that are big enough to contain meaningful text
                        if img.width > 100 and img.height > 100:
                            images.append(img)
                    except Exception as e:
                        logger.error(f"Error extracting image from PowerPoint slide: {e}")
            except Exception as e:
                logger.error(f"Error processing shape in PowerPoint slide: {e}")
        
        return images

    async def _extract_images_with_positions(self, slide) -> List[Tuple[Image.Image, float, float]]:
        """
        Extract images from a PowerPoint slide along with their positions.
        
        Args:
            slide: python-pptx slide object
            
        Returns:
            List[Tuple[Image.Image, float, float]]: List of tuples containing (image, left, top)
        """
        images_with_pos = []
        
        for shape in slide.shapes:
            try:
                # Check if shape is a picture
                if hasattr(shape, 'image'):
                    try:
                        # Extract the image data
                        image_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        
                        # Only process images that are big enough to contain meaningful text
                        if img.width > 100 and img.height > 100:
                            # Get position information
                            left = shape.left
                            top = shape.top
                            images_with_pos.append((img, left, top))
                    except Exception as e:
                        logger.error(f"Error extracting image from PowerPoint slide: {e}")
            except Exception as e:
                logger.error(f"Error processing shape in PowerPoint slide: {e}")
        
        # Sort by vertical position (top to bottom)
        images_with_pos.sort(key=lambda x: x[2])
        
        return images_with_pos 